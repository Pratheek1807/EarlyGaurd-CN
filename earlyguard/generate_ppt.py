"""
EarlyGuard — Demo PPT Generator  (light theme v3)
3 slides: Problem + Dataset | Pipeline + 10 Signals | Performance + Collections + ROI
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours (light theme) ─────────────────────────────────────
BG       = RGBColor(0xF8, 0xFA, 0xFF)
SURFACE  = RGBColor(0xFF, 0xFF, 0xFF)
S2       = RGBColor(0xF1, 0xF5, 0xF9)
S3       = RGBColor(0xE2, 0xE8, 0xF0)
DARK     = RGBColor(0x0F, 0x17, 0x2A)

ACCENT   = RGBColor(0x1D, 0x4E, 0xD8)
GREEN    = RGBColor(0x15, 0x80, 0x3D)
RED      = RGBColor(0xB9, 0x1C, 0x1C)
ORANGE   = RGBColor(0xC2, 0x41, 0x0C)
YELLOW   = RGBColor(0xA1, 0x62, 0x07)
PURPLE   = RGBColor(0x6D, 0x28, 0xD9)

ACCENT_L = RGBColor(0xDB, 0xEA, 0xFE)
GREEN_L  = RGBColor(0xDC, 0xFC, 0xE7)
RED_L    = RGBColor(0xFE, 0xE2, 0xE2)
ORANGE_L = RGBColor(0xFF, 0xED, 0xD5)
PURPLE_L = RGBColor(0xED, 0xE9, 0xFE)
YELLOW_L = RGBColor(0xFE, 0xF9, 0xC3)

TEXT     = RGBColor(0x1E, 0x29, 0x3B)
MUTED    = RGBColor(0x47, 0x55, 0x69)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_SUB = RGBColor(0x94, 0xA3, 0xB8)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────
def tb(slide, x, y, w, h, text, size=12, bold=False, color=TEXT,
       align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def rect(slide, x, y, w, h, fill=SURFACE, border=None, border_w=Pt(0.75)):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape


def add_bg(slide):
    r = rect(slide, 0, 0, 13.33, 7.5, fill=BG)
    r.line.fill.background()


def slide_header(slide, num):
    rect(slide, 0, 0, 13.33, 0.88, fill=DARK)
    tb(slide, 0.4, 0.08, 3.5, 0.46, "EarlyGuard", size=26, bold=True, color=WHITE)
    tb(slide, 0.4, 0.55, 8, 0.28, "ML-Powered Early Warning for Loan Defaults", size=10, color=HEADER_SUB)
    tb(slide, 11.4, 0.32, 1.8, 0.26, f"SLIDE {num} / 3", size=9, color=HEADER_SUB, align=PP_ALIGN.RIGHT)
    rect(slide, 0, 0.88, 13.33, 0.045, fill=ACCENT)


def slide_footer(slide):
    rect(slide, 0, 7.2, 13.33, 0.3, fill=S2)
    tb(slide, 0.4, 7.23, 13, 0.24, "Credit Nirvana  ·  EarlyGuard Demo  ·  Confidential", size=9, color=MUTED)


def slide_title(slide, text):
    tb(slide, 0.5, 1.05, 12.5, 0.52, text, size=23, bold=True, color=TEXT)
    rect(slide, 0.5, 1.62, 0.5, 0.05, fill=ACCENT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 1 — The Problem & What We Built
# ══════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
add_bg(s1)
slide_header(s1, 1)
slide_footer(s1)
slide_title(s1, "The Problem — and What We Built to Solve It")

# ── LEFT: Problem card ─────────────────────────────────────────
rect(s1, 0.4, 1.75, 5.9, 5.32, fill=RED_L, border=RED, border_w=Pt(1.2))
rect(s1, 0.4, 1.75, 5.9, 0.07, fill=RED)
tb(s1, 0.62, 1.87, 5.5, 0.28, "THE PROBLEM", size=9, bold=True, color=RED)

tb(s1, 0.62, 2.2, 5.4, 0.88,
   "Lenders discover defaults only after they happen. By then the account is already "
   "delinquent and the only option is expensive field recovery.",
   size=12, color=MUTED)

problems = [
    ("No early signal",      "Decisions based on lagging DPD data only — no predictive view"),
    ("One-size response",    "All at-risk accounts treated the same — zero risk stratification"),
    ("Ops cost explosion",   "Post-default: ₹101.50/account  (field agent ₹100 + call ₹1 + SMS ₹0.50)"),
    ("Loan loss",            "50% of at-risk accounts default without any proactive intervention"),
]
for i, (title, desc) in enumerate(problems):
    ry = 3.2 + i * 0.73
    rect(s1, 0.58, ry, 5.55, 0.64, fill=SURFACE, border=RED, border_w=Pt(0.5))
    tb(s1, 0.78, ry+0.07, 2.6, 0.25, title, size=11, bold=True, color=RED)
    tb(s1, 0.78, ry+0.35, 5.1, 0.24, desc,  size=10, color=MUTED)

# ── RIGHT top: Solution card ────────────────────────────────────
rect(s1, 6.6, 1.75, 6.3, 2.35, fill=ACCENT_L, border=ACCENT, border_w=Pt(1.2))
rect(s1, 6.6, 1.75, 6.3, 0.07, fill=ACCENT)
tb(s1, 6.82, 1.87, 6.0, 0.28, "THE SOLUTION — EARLYGUARD", size=9, bold=True, color=ACCENT)
tb(s1, 6.82, 2.18, 6.0, 1.08,
   "Score every loan account before it defaults — 30 days and 60 days ahead.\n"
   "Route only genuinely at-risk accounts to collections with a precise action, timing, tone, and offer.",
   size=12, color=TEXT)
tb(s1, 6.82, 3.14, 6.0, 0.4,
   "Result: catch slow-burn accounts that a 30d-only system misses entirely.",
   size=11, italic=True, color=ACCENT)

# ── RIGHT bottom: Dataset card ──────────────────────────────────
rect(s1, 6.6, 4.25, 6.3, 2.82, fill=GREEN_L, border=GREEN, border_w=Pt(1.2))
rect(s1, 6.6, 4.25, 6.3, 0.07, fill=GREEN)
tb(s1, 6.82, 4.35, 6.0, 0.28, "DATASET — 5,000 LOAN ACCOUNTS  ·  6 DATA SOURCES", size=9, bold=True, color=GREEN)

datasets = [
    ("🏦", "Loan Accounts",     "EMI, outstanding balance, employment type, product"),
    ("🏧", "Bank Statements",   "Monthly inflow, surplus, expense ratios — last 6 months"),
    ("📊", "Bureau Data",       "Credit score history, hard enquiries, missed payments"),
    ("📞", "Telecall History",  "Response rates, PTP outcomes, sentiment, contact avoidance"),
    ("📋", "EPFO / Employment", "Contribution gaps, job changes, GST & ITR filing status"),
    ("💳", "Payment History",   "EMI delays, partial payments, missed payments — 6 months"),
]
for i, (icon, name, desc) in enumerate(datasets):
    ry = 4.78 + i * 0.355
    tb(s1, 6.82, ry, 0.32, 0.3, icon, size=10, color=TEXT)
    tb(s1, 7.2,  ry, 1.85, 0.3, name, size=10.5, bold=True, color=TEXT)
    tb(s1, 9.1,  ry, 3.65, 0.3, desc, size=9.5,  color=MUTED)

# Assumption notes
rect(s1, 6.6, 6.9, 6.3, 0.28, fill=YELLOW_L)
tb(s1, 6.75, 6.91, 6.0, 0.12,
   "✦  Works with any combination of sources — even loan data alone is enough to score accounts.",
   size=7.5, italic=True, color=YELLOW)
tb(s1, 6.75, 7.04, 6.0, 0.12,
   "✦  Bank statements assumed pre-processed: categorised transactions (salary credit, UPI inflows, business turnover, rent, etc.) already extracted by a statement analyser.",
   size=7.5, italic=True, color=YELLOW)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 2 — Pipeline + 10 Engineered Signals
# ══════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_bg(s2)
slide_header(s2, 2)
slide_footer(s2)
slide_title(s2, "How It Works — 4 Stages, 10 Signals, One Risk Score")

# ── 4-stage pipeline strip ─────────────────────────────────────
stages = [
    (ACCENT,  "STAGE 1", "6 Raw Datasets",    "Loan · Bank · Bureau\nTelecall · EPFO · Payment"),
    (PURPLE,  "STAGE 2", "10 Signals",        "Domain-scored 0–1\nper account"),
    (ORANGE,  "STAGE 3", "Stacked Ensemble",  "LR + RF + XGBoost\n→ Meta LR combiner"),
    (GREEN,   "STAGE 4", "Risk Score",        "Tier · Action · Timing\nTone · Offer"),
]
box_w  = 2.88
box_xs = [0.4, 3.57, 6.74, 9.91]

for i, (color, label, title, desc) in enumerate(stages):
    x = box_xs[i]
    rect(s2, x, 1.72, box_w, 1.92, fill=SURFACE, border=color, border_w=Pt(1.2))
    rect(s2, x, 1.72, box_w, 0.07, fill=color)
    tb(s2, x+0.14, 1.83, box_w-0.22, 0.24, label, size=8.5, bold=True, color=color)
    tb(s2, x+0.14, 2.1,  box_w-0.22, 0.5,  title, size=15, bold=True, color=TEXT)
    tb(s2, x+0.14, 2.65, box_w-0.22, 0.9,  desc,  size=10, color=MUTED)
    if i < 3:
        tb(s2, x+box_w+0.05, 2.52, 0.38, 0.35, "→", size=20, color=S3, align=PP_ALIGN.CENTER)

# ── Signals header ─────────────────────────────────────────────
rect(s2, 0.4, 3.75, 12.5, 0.3, fill=S2)
tb(s2, 0.55, 3.79, 12.2, 0.22,
   "10 ENGINEERED SIGNALS  ·  Scored 0 (healthy) → 1 (stressed)  ·  Signals & weights fully configurable — add, remove, or retrain per client",
   size=9, bold=True, color=MUTED)

# ── Signal cards: 2 columns × 5 rows ──────────────────────────
# Each card: name (bold) | one-liner (italic) | sub-features (right half)
# Card height = 0.60", card width = 6.25"

signals = [
    # (Name, one-liner, [sub-features with weights], color)
    (
        "Payment Deterioration",
        "Is the customer short-paying or skipping EMIs right now?",
        [
            ("Partial payment ratio",        "40%"),
            ("Count of missed EMIs (3m)",    "35%"),
            ("Average days late (3m)",       "25%"),
        ],
        ORANGE,
    ),
    (
        "Financial Cushion Trend",
        "Is the bank balance shrinking month after month?",
        [
            ("Net surplus slope (3m)",       "45%"),
            ("Expense-to-income ratio",      "30%"),
            ("Days balance near zero",       "25%"),
        ],
        ACCENT,
    ),
    (
        "Income Reliability",
        "Is the customer's income stable and arriving on time?",
        [
            ("Salaried — salary delay days", "40%"),
            ("EPFO zero-contribution months","35%"),
            ("Job change detected",          "25%"),
            ("Self-emp — GST filing stress", "40%"),
            ("UPI inflow trend declining",   "35%"),
            ("ITR not filed flag",           "25%"),
        ],
        PURPLE,
    ),
    (
        "Bureau Stress",
        "Is the customer's credit score dropping and missing elsewhere?",
        [
            ("Credit score drop (vs 6m ago)","40%"),
            ("Missed payments at other lenders","35%"),
            ("Hard enquiries last 30 days",  "25%"),
        ],
        RED,
    ),
    (
        "Contact Avoidance",
        "Is the customer picking up calls or going completely silent?",
        [
            ("Response rate (30d)",          "35%"),
            ("Consecutive unanswered calls", "35%"),
            ("Days since last contact",      "30%"),
        ],
        ORANGE,
    ),
    (
        "Debt Pressure Ratio",
        "How much of the income goes to paying all EMIs combined?",
        [
            ("Our EMI + other lender EMIs",  "—"),
            ("÷ total monthly inflow",       "direct ratio, uncapped"),
        ],
        RED,
    ),
    (
        "PTP Reliability",
        "When the customer promises to pay, do they follow through?",
        [
            ("Broken PTPs last 3 months",    "50%"),
            ("Historical fulfilment rate",   "50%"),
        ],
        YELLOW,
    ),
    (
        "Sentiment Distress",
        "Does the customer sound stressed or hostile on calls?",
        [
            ("Sentiment score (last call)",  "35%"),
            ("Dominant emotion type",        "30%"),
            ("Hardship flag mentioned",      "20%"),
            ("3-month sentiment trend",      "15%"),
        ],
        PURPLE,
    ),
    (
        "Recency-Weighted Stress",
        "Recent missed payments count far more — catches sudden drops.",
        [
            ("6-month payment history",      ""),
            ("Weights oldest→newest:",       "5%→8%→12%→17%→25%→33%"),
            ("Newest payment = 6× oldest",   ""),
        ],
        ORANGE,
    ),
    (
        "Employment Event Recency",
        "Did the customer recently lose a job or show business trouble?",
        [
            ("Salaried — job change recency","primary"),
            ("EPFO zero-contribution months","secondary"),
            ("Self-emp — GST trend decline", "primary"),
            ("ITR not filed (self-emp)",     "secondary"),
        ],
        ACCENT,
    ),
]

CARD_H   = 0.625
CARD_W   = 6.25
NAME_W   = 2.55    # left portion for name + one-liner
FEAT_X   = 2.78    # x offset for sub-features (relative to card x)
FEAT_W   = 3.35    # sub-features column width
cols_x   = [0.4, 6.84]

for i, (name, oneliner, subfeats, color) in enumerate(signals):
    col = i % 2
    row = i // 2
    cx  = cols_x[col]
    cy  = 4.14 + row * CARD_H

    # Card background + left colour bar
    rect(s2, cx, cy, CARD_W, CARD_H - 0.04, fill=SURFACE, border=S3, border_w=Pt(0.5))
    rect(s2, cx, cy, 0.06, CARD_H - 0.04, fill=color)

    # Name + one-liner (left half)
    tb(s2, cx+0.13, cy+0.04,  NAME_W, 0.24, name,     size=10, bold=True, color=TEXT)
    tb(s2, cx+0.13, cy+0.30,  NAME_W, 0.26, oneliner, size=8,  italic=True, color=MUTED)

    # Sub-features (right half) — compact, one line each
    feat_lines = "  ·  ".join(
        f"{lbl} ({wt})" if wt and wt not in ("", "—", "primary", "secondary")
        else f"{lbl}{(' ['+wt+']') if wt and wt not in ('', '—') else ''}"
        for lbl, wt in subfeats
    )
    tb(s2, cx+FEAT_X, cy+0.04, FEAT_W, CARD_H-0.1, feat_lines, size=8, color=color)

# Customisation banner
rect(s2, 0.4, 7.02, 12.5, 0.16, fill=ACCENT_L)
tb(s2, 0.55, 7.03, 12.2, 0.14,
   "Signals can be added or removed  ·  Sub-weights adjusted by the credit team without retraining  ·  New weights take effect automatically on next run.",
   size=8.5, bold=True, color=ACCENT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 3 — Performance + What Collections Gets + Business ROI
# ══════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_bg(s3)
slide_header(s3, 3)
slide_footer(s3)
slide_title(s3, "Validated Performance + What Collections Gets + Business ROI")

# ── TOP: Model metrics strip ────────────────────────────────────
rect(s3, 0.4, 1.72, 12.5, 1.52, fill=SURFACE, border=ACCENT, border_w=Pt(1.2))
rect(s3, 0.4, 1.72, 12.5, 0.07, fill=ACCENT)
tb(s3, 0.62, 1.82, 10, 0.26,
   "MODEL VALIDATION — HELD-OUT 30% TEST SET  ·  STACKED ENSEMBLE: LR + RF + XGBoost → Meta LR Combiner",
   size=9, bold=True, color=ACCENT)

metrics = [
    ("AUC-ROC  (30-day model)", "0.8596", "target ≥ 0.72", "PASS ✓"),
    ("AUC-ROC  (60-day model)", "0.8736", "target ≥ 0.68", "PASS ✓"),
    ("False Positive Rate",     "12.51%", "target ≤ 25%",  "PASS ✓"),
]
for i, (label, val, target, verdict) in enumerate(metrics):
    x = 0.55 + i * 4.2
    rect(s3, x, 2.12, 4.0, 1.02, fill=S2)
    tb(s3, x+0.14, 2.18, 3.7, 0.26, label,   size=10, color=MUTED)
    tb(s3, x+0.14, 2.44, 2.0, 0.46, val,     size=24, bold=True, color=GREEN)
    tb(s3, x+0.14, 2.92, 2.2, 0.18, target,  size=8.5, color=MUTED)
    tb(s3, x+2.6,  2.92, 1.3, 0.18, verdict, size=9,  bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

# ── HERO: Current output — ACTION ────────────────────────────
rect(s3, 0.4, 3.36, 7.6, 2.25, fill=ACCENT_L, border=ACCENT, border_w=Pt(2.0))
rect(s3, 0.4, 3.36, 7.6, 0.08, fill=ACCENT)
tb(s3, 0.62, 3.48, 7.2, 0.3,  "WHAT EARLYGUARD OUTPUTS TODAY — RECOMMENDED ACTION", size=10, bold=True, color=ACCENT)
tb(s3, 0.62, 3.82, 7.2, 0.22,
   "Every flagged account gets one specific action driven by its highest-stress signal.",
   size=10, italic=True, color=MUTED)

actions = [
    "EMI Relief Call",  "Hardship Support Call",  "PTP Re-engagement",
    "Empathy Outreach", "Bureau Advisory",         "Urgent Follow-up Call",
    "Restructure Discussion", "Standard Follow-up", "No Intervention (Low risk)",
]
# 3 columns of 3
for i, act in enumerate(actions):
    col, row = i % 3, i // 3
    ax = 0.58 + col * 2.5
    ay = 4.1  + row * 0.38
    rect(s3, ax, ay, 2.35, 0.3, fill=SURFACE, border=ACCENT, border_w=Pt(0.5))
    tb(s3, ax+0.1, ay+0.05, 2.2, 0.22, act, size=9, bold=True, color=ACCENT)

# ── Business ROI (right of action hero) ───────────────────────
rect(s3, 8.28, 3.36, 4.62, 2.25, fill=GREEN_L, border=GREEN, border_w=Pt(1.5))
rect(s3, 8.28, 3.36, 4.62, 0.08, fill=GREEN)
tb(s3, 8.45, 3.48, 4.3, 0.28, "BUSINESS CASE — 10 ACCOUNTS", size=10, bold=True, color=GREEN)

roi_rows = [
    ("At-risk accounts flagged",           "10"),
    ("Baseline default rate",              "50%  →  5 default"),
    ("Loan amount at risk (₹1L each)",    "₹5,00,000"),
    ("EarlyGuard recovery rate",           "40%  →  4 saved"),
    ("Post-default ops cost / account",   "₹101.50"),
    ("EarlyGuard early ops cost",          "₹1.50 / account"),
    ("Loan amount protected",              "₹4,00,000"),
    ("Net ROI on at-risk portfolio",       "~80%"),
]
for i, (lbl, val) in enumerate(roi_rows):
    ry = 3.82 + i * 0.22
    tb(s3, 8.42, ry, 3.0, 0.2, lbl, size=8.5, color=MUTED)
    tb(s3, 11.5, ry, 1.3, 0.2, val, size=8.5, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

# ── Future Scope strip ─────────────────────────────────────────
rect(s3, 0.4, 5.72, 12.5, 1.42, fill=S2, border=S3, border_w=Pt(0.75))
rect(s3, 0.4, 5.72, 12.5, 0.07, fill=ORANGE)
tb(s3, 0.62, 5.82, 12.0, 0.26, "FUTURE SCOPE — TIMING, TONE & OFFER", size=10, bold=True, color=ORANGE)
tb(s3, 0.62, 6.1,  12.0, 0.3,
   "These three outputs are an already-solved problem for most collections teams. "
   "EarlyGuard's ACTION + risk score can plug directly into existing business-rule engines or CRM systems that the client already operates.",
   size=9.5, italic=True, color=MUTED)

future_cards = [
    (GREEN,  "TIMING",
     "Best contact window — fetched from client's CRM / IVR system using historical response data per account"),
    (ORANGE, "TONE",
     "Empathy / vulnerability flag — sourced from existing sentiment tools or collections playbook rules"),
    (PURPLE, "OFFER",
     "Restructure / holiday / partial — pulled from credit team's business-rule engine against account profile"),
]
for i, (color, label, desc) in enumerate(future_cards):
    fx = 0.5 + i * 4.15
    rect(s3, fx, 6.44, 3.95, 0.65, fill=SURFACE, border=color, border_w=Pt(0.75))
    tb(s3, fx+0.12, 6.48, 1.0,  0.22, label, size=9,   bold=True, color=color)
    tb(s3, fx+1.15, 6.48, 2.72, 0.56, desc,  size=8.5, color=MUTED)


# ── Save ──────────────────────────────────────────────────────
out = "/home/rohitsaini/Documents/perfios/code/EarlyGaurd-CN/earlyguard/EarlyGuard_Demo.pptx"
prs.save(out)
print(f"Saved → {out}")
